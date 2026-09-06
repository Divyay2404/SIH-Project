"""Extract export-ready text chunks from PDF and PowerPoint curriculum material."""

import os
import re
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any, Dict, List

from app.ingestion.pdf_parser import pdf_parser_engine

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    Presentation = None
    PPTX_AVAILABLE = False


class UnsupportedDocumentError(ValueError):
    """Raised when a document format cannot be extracted in the current runtime."""


class CurriculumDocumentParser:
    SUPPORTED_EXTENSIONS = {".pdf", ".ppt", ".pptx"}

    def _slide_text(self, slide: Any) -> List[str]:
        """Collect visible text from text boxes, grouped shapes, and tables."""
        extracted: List[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    extracted.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        extracted.append(" | ".join(cells))
            if getattr(shape, "shape_type", None) == 6:  # GROUP
                extracted.extend(self._shape_text(shape.shapes))
        return extracted

    @staticmethod
    def _shape_text(shapes: Any) -> List[str]:
        extracted: List[str] = []
        for shape in shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                extracted.append(shape.text.strip())
        return extracted

    def _parse_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        if not PPTX_AVAILABLE:
            raise UnsupportedDocumentError("PowerPoint parsing requires the python-pptx package.")
        try:
            presentation = Presentation(file_path)
        except Exception as error:
            raise UnsupportedDocumentError(f"Unable to open this PowerPoint file: {error}") from error

        chunks: List[Dict[str, Any]] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            text = "\n".join(self._slide_text(slide)).strip()
            if text:
                chunks.append({
                    "page": slide_number,
                    "text": text,
                    "bbox": [0.0, 0.0, 1.0, 1.0],
                    "block_type": "slide",
                })
        return chunks

    def _parse_legacy_ppt(self, file_path: str) -> List[Dict[str, Any]]:
        """Convert legacy .ppt files to .pptx when LibreOffice is available."""
        converter = os.environ.get("LIBREOFFICE_PATH") or shutil.which("soffice") or shutil.which("libreoffice")
        if not converter:
            return self._extract_legacy_ppt_text(file_path)
        with tempfile.TemporaryDirectory() as output_dir:
            result = subprocess.run(
                [converter, "--headless", "--convert-to", "pptx", "--outdir", output_dir, file_path],
                capture_output=True,
                text=True,
                timeout=60,
                check=False,
            )
            converted = Path(output_dir, f"{Path(file_path).stem}.pptx")
            if result.returncode != 0 or not converted.exists():
                raise UnsupportedDocumentError("The legacy .ppt file could not be converted to a readable PowerPoint file.")
            return self._parse_pptx(str(converted))

    @staticmethod
    def _extract_legacy_ppt_text(file_path: str) -> List[Dict[str, Any]]:
        """Recover embedded Unicode slide text from a legacy binary PowerPoint file.

        LibreOffice produces richer results when it is available. This fallback
        keeps legacy decks usable in serverless environments where conversion
        software cannot be installed.
        """
        content = Path(file_path).read_bytes()
        candidates = re.findall(rb"(?:[\x20-\x7e]\x00){3,}", content)
        text_parts: List[str] = []
        seen = set()
        for candidate in candidates:
            value = candidate.decode("utf-16-le", errors="ignore").strip()
            if value and value not in seen:
                seen.add(value)
                text_parts.append(value)
        if not text_parts:
            raise UnsupportedDocumentError(
                "No readable text was found in this legacy .ppt file. Try saving it as .pptx and uploading again."
            )
        return [{
            "page": 1,
            "text": "\n".join(text_parts),
            "bbox": [0.0, 0.0, 1.0, 1.0],
            "block_type": "legacy_slide",
        }]

    def parse(self, file_path: str) -> List[Dict[str, Any]]:
        extension = Path(file_path).suffix.lower()
        if extension == ".pdf":
            return pdf_parser_engine.parse_pdf(file_path)
        if extension == ".pptx":
            return self._parse_pptx(file_path)
        if extension == ".ppt":
            return self._parse_legacy_ppt(file_path)
        raise UnsupportedDocumentError("Supported formats are PDF, PPTX, and PPT.")


document_parser = CurriculumDocumentParser()
