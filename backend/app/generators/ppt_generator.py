"""Generate editable lecture decks from the chunks extracted from an uploaded document."""

import io
from typing import Any, Dict, List

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class PresentationGenerator:
    def __init__(self):
        self.available = PPTX_AVAILABLE

    @staticmethod
    def _points(text: str, limit: int = 5) -> List[str]:
        """Turn a text block into concise, source-preserving slide points."""
        lines = [line.strip(" -•\t") for line in text.replace("\r", "").split("\n")]
        points = [line for line in lines if line]
        if len(points) == 1:
            sentences = [sentence.strip() for sentence in text.replace("\n", " ").split(".")]
            points = [sentence for sentence in sentences if sentence]
        return points[:limit] or ["No extractable content was found for this section."]

    def generate_ppt_deck(self, document: Dict[str, Any]) -> bytes:
        """Build a deck whose slide text and notes come from one uploaded document."""
        if not self.available:
            raise RuntimeError("python-pptx is required to generate PowerPoint files.")

        title = document["title"]
        chunks = document.get("chunks", [])
        if not chunks:
            raise ValueError("The selected document has no extracted content to export.")

        prs = Presentation()
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = title
        title_slide.placeholders[1].text = "Lecture deck generated from the uploaded curriculum material"
        title_slide.notes_slide.notes_text_frame.text = f"Introduce {title} using the uploaded document as the source of record."

        content_layout = prs.slide_layouts[1]
        for index, chunk in enumerate(chunks[:6], start=1):
            slide = prs.slides.add_slide(content_layout)
            lines = self._points(chunk["text"])
            slide.shapes.title.text = lines[0][:90] if lines else f"{title}: Section {index}"
            text_frame = slide.shapes.placeholders[1].text_frame
            text_frame.clear()
            for point_index, point in enumerate(lines[1:] or lines[:1]):
                paragraph = text_frame.paragraphs[0] if point_index == 0 else text_frame.add_paragraph()
                paragraph.text = point[:280]
                paragraph.level = 0
            slide.notes_slide.notes_text_frame.text = (
                f"Source page {chunk.get('page', 'unknown')}. Explain the displayed points "
                "and refer back to the uploaded material for examples and detail."
            )

        summary = prs.slides.add_slide(content_layout)
        summary.shapes.title.text = f"{title}: Key Takeaways"
        summary_frame = summary.shapes.placeholders[1].text_frame
        summary_frame.clear()
        for index, chunk in enumerate(chunks[:4]):
            paragraph = summary_frame.paragraphs[0] if index == 0 else summary_frame.add_paragraph()
            paragraph.text = self._points(chunk["text"], limit=1)[0][:280]
        summary.notes_slide.notes_text_frame.text = "Review the main ideas drawn from the uploaded document."

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()


ppt_generator = PresentationGenerator()
