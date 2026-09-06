"""Regression tests that prevent educator exports from reverting to BST placeholders."""

import io
import unittest

from app.generators.pdf_generator import pdf_generator
from app.generators.ppt_generator import ppt_generator


DOCUMENT = {
    "title": "Operating Systems Scheduling",
    "chunks": [
        {"page": 1, "text": "Process Scheduling\nThe scheduler selects a ready process for CPU execution."},
        {"page": 2, "text": "Context Switching\nThe operating system saves one process state before restoring another."},
    ],
}


class TestDynamicEducatorExports(unittest.TestCase):
    def test_presentation_uses_uploaded_document_content(self):
        if not ppt_generator.available:
            self.skipTest("python-pptx is not installed")
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(ppt_generator.generate_ppt_deck(DOCUMENT)))
        text = "\n".join(shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text"))
        self.assertIn("Operating Systems Scheduling", text)
        self.assertIn("Process Scheduling", text)
        self.assertNotIn("Binary Search Trees", text)

    def test_handout_uses_uploaded_document_content(self):
        if not pdf_generator.available:
            self.skipTest("ReportLab is not installed")
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(pdf_generator.generate_handout_pdf(DOCUMENT)))
        text = "\n".join(page.extract_text() or "" for page in reader.pages)
        self.assertIn("Operating Systems Scheduling", text)
        self.assertIn("Context Switching", text)
        self.assertNotIn("Binary Search Trees", text)


if __name__ == "__main__":
    unittest.main()
